from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).resolve().parent / "Sign_Language_Project_15_Slides.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], level2: dict[int, list[str]] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        if level2 and i in level2:
            for sub in level2[i]:
                sp = body.add_paragraph()
                sp.text = sub
                sp.level = 1
                sp.font.size = Pt(18)


def add_metrics_table_slide(prs: Presentation, title: str, columns: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    n_rows = len(rows) + 1
    n_cols = len(columns)
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2)
    ).table

    for j, col in enumerate(columns):
        cell = table.cell(0, j)
        cell.text = col
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(16)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(15)


def add_final_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Thank You"
    left, top, width, height = Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = "Questions and Discussion"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True


def build() -> Path:
    prs = Presentation()

    add_title_slide(
        prs,
        "Sign Language Detection System",
        "ML 257 Project\nLetters + Words + Decoder",
    )

    add_bullets_slide(
        prs,
        "1. Problem Statement and Motivation",
        [
            "Enable accessible communication by converting signs to text.",
            "Target both fingerspelling (letters) and vocabulary signs (words).",
            "Build a practical system: trainable, evaluable, and demo-ready.",
        ],
    )

    add_bullets_slide(
        prs,
        "2. Project Objectives",
        [
            "Compare multiple ML paradigms for sign recognition.",
            "Deliver robust inference with post-processing for readability.",
            "Provide reproducible scripts, metrics, and deployment options.",
        ],
    )

    add_bullets_slide(
        prs,
        "3. System Architecture",
        [
            "Part 1: Letter recognition (classical ML + deep vision + YOLO).",
            "Part 2: Word recognition (BiLSTM and Transformer).",
            "Part 3: Decoder for temporal smoothing and sentence assembly.",
            "Flask UI for upload/live camera and end-to-end demonstrations.",
        ],
    )

    add_bullets_slide(
        prs,
        "4. Dataset and Splits",
        [
            "Part 1 landmarks: 2424 total (Train 1696, Val 364, Test 364).",
            "Part 1 images: 2515 total (Train 1759, Val 378, Test 378).",
            "YOLO dataset: Train 2012, Val 503.",
            "Part 2 WLASL sequences: 568 total (Train 407, Val 100, Test 61).",
        ],
    )

    add_bullets_slide(
        prs,
        "5. Part 1 Preprocessing",
        [
            "MediaPipe extracts 21 hand landmarks for classical models.",
            "Feature vectors saved as X.npy, labels as y.npy, mapping in label_map.npy.",
            "Deep models load raw RGB images resized to 64x64 and normalized.",
        ],
    )

    add_bullets_slide(
        prs,
        "6. Part 1 Models",
        [
            "Classical: SVM (RBF), Random Forest, MLP.",
            "Deep: Custom CNN, MobileNetV2, ResNet-18, VGG-11-BN.",
            "YOLO classification branch evaluated on val split.",
        ],
        level2={
            1: [
                "Transfer models use ImageNet-pretrained backbones by default.",
                "Early stopping and LR scheduling improve convergence.",
            ]
        },
    )

    add_metrics_table_slide(
        prs,
        "7. Part 1 Results (Letters)",
        ["Model", "Accuracy", "Macro F1"],
        [
            ["Random Forest", "0.9066", "0.9027"],
            ["SVM", "0.9011", "0.8977"],
            ["MLP", "0.8242", "0.8173"],
            ["CNN", "0.9392", "0.9384"],
            ["MobileNetV2", "0.9868", "0.9867"],
            ["ResNet-18", "0.9894", "0.9892"],
            ["VGG-11-BN", "0.9815", "0.9816"],
            ["YOLO (val)", "0.9622", "0.9617"],
        ],
    )

    add_bullets_slide(
        prs,
        "8. Part 2 Preprocessing and Features",
        [
            "WLASL videos converted to fixed 30-frame sequences.",
            "Per-frame feature size: 225 (left hand + right hand + pose).",
            "Training augmentation: noise, frame dropout, temporal scaling.",
            "Weighted sampler handles class imbalance.",
        ],
    )

    add_bullets_slide(
        prs,
        "9. Part 2 Models",
        [
            "BiLSTM: bidirectional sequence encoder over 30 frames.",
            "Transformer: multi-head self-attention with positional encoding.",
            "Both optimized with Adam, LR scheduling, early stopping.",
        ],
    )

    add_metrics_table_slide(
        prs,
        "10. Part 2 Results (Words, 100 classes)",
        ["Model", "Top-1", "Top-5", "Best Val Top-1"],
        [
            ["Transformer", "34.4%", "70.5%", "41.0%"],
            ["BiLSTM", "24.6%", "59.0%", "28.0%"],
        ],
    )

    add_bullets_slide(
        prs,
        "11. Part 3 Decoder",
        [
            "Input: frame-wise letter predictions plus confidence scores.",
            "Temporal smoothing reduces flicker and unstable transitions.",
            "Lexical decoding corrects noisy letter sequences to plausible words.",
            "Beam-search sentence builder outputs best hypotheses.",
        ],
    )

    add_bullets_slide(
        prs,
        "12. Evaluation Strategy",
        [
            "Use held-out test sets for unbiased final reporting.",
            "Track accuracy, macro precision/recall/F1, confusion matrices.",
            "For word models, include Top-1 and Top-5 accuracy.",
        ],
        level2={
            2: [
                "Top-1: first guess must match.",
                "Top-5: correct label appears in top 5 guesses.",
            ]
        },
    )

    add_bullets_slide(
        prs,
        "13. Key Findings",
        [
            "Deep image models outperform landmark-based classical models for letters.",
            "ResNet-18 is the strongest Part 1 model on current test split.",
            "Transformer outperforms BiLSTM on WLASL word recognition.",
            "Decoder is crucial for turning noisy frame predictions into usable text.",
        ],
    )

    add_bullets_slide(
        prs,
        "14. Limitations and Future Work",
        [
            "Word-level performance limited by small/imbalanced Part 2 data.",
            "Dead/invalid WLASL source links reduce usable samples.",
            "Future: larger balanced datasets, multimodal fusion, language-model integration.",
            "Future: optimize real-time inference and mobile deployment.",
        ],
    )

    add_final_slide(prs)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build()
    print(f"Saved: {out}")
